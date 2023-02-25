import PyPDF2
import pptx
from pptx.util import Inches
from transformers import BertTokenizer, BertForSequenceClassification

# Load the pre-trained BERT tokenizer and model
tokenizer = BertTokenizer.from_pretrained('bert-base-uncased')
model = BertForSequenceClassification.from_pretrained('bert-base-uncased')

# Open the PDF file and read its contents
pdf_file = open('research_paper.pdf', 'rb')
pdf_reader = PyPDF2.PdfFileReader(pdf_file)
text = ''
for page_num in range(pdf_reader.getNumPages()):
    page = pdf_reader.getPage(page_num)
    text += page.extractText()

# Split the text into sentences
sentences = text.split('.')

# Use BERT to predict the label for each sentence
labels = []
for i in range(len(sentences)):
    inputs = tokenizer(
        sentences[i], return_tensors='pt', truncation=True, padding=True)
    outputs = model(inputs['input_ids'],
                    attention_mask=inputs['attention_mask'])
    predicted_label = outputs[0].argmax(dim=1).item()
    labels.append(predicted_label)

# Identify the key insights and findings
key_sentences = []
for i in range(len(sentences)):
    if labels[i] == 1:
        key_sentences.append(sentences[i])

# Create a PowerPoint presentation with the key insights and findings
ppt = pptx.Presentation()
for i in range(len(key_sentences)):
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"Key Insight {i+1}"
    body = slide.placeholders[1]
    body.text = key_sentences[i]
    img = slide.shapes.add_picture(
        'image.jpg', Inches(4), Inches(4), height=Inches(2))
ppt.save('research_presentation.pptx')
