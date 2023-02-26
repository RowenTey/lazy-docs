import os
import pptx
from pptx.util import Pt, Inches


def generate_ppt(content, ppt_title="Research Paper Summary", filename="research_summary.ppt"):
    directory = os.path.join(os.getcwd(), '..', 'web-app/public')
    if not os.path.exists(directory):
        os.makedirs(directory)
    filepath = os.path.join(directory, filename)

    # Create a PowerPoint presentation with the key insights and findings
    ppt = pptx.Presentation()

    # add title page
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    title.text = ppt_title
    title.left = Inches(0.5)  # adjust the left position as necessary
    title.top = Inches(3)
    title.width = Inches(9)
    title.height = Inches(1)
    title.horz_cent = True
    title.vert_cent = True

    # remove body placeholder for title page
    body_shape = slide.placeholders[1]
    body_shape.text = ""
    body_shape.element.clear()

    # add insights
    for section, insights in content.items():
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])

        title = slide.shapes.title
        title.text = section.upper()

        body = slide.placeholders[1]
        body.text = insights

        # Set the font size of the body text to 16
        for para in body.text_frame.paragraphs:
            para.font.size = Pt(17)
    ppt.save(filepath)


if __name__ == "__main__":
    tmp = ["hello\nhaha", "world"]
    generate_ppt(tmp, "test1.ppt")
